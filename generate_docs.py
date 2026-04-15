from docx import Document
from pptx import Presentation
import os

# 1. Generate Word Document
doc = Document()
doc.add_heading('Sunset Agent OS: Architectural Documentation', 0)
doc.add_paragraph('Authors: Arshad Pasha and Arun')

doc.add_heading('1. Overview & Core Vision', level=1)
doc.add_paragraph(
    'Sunset Agent OS is a paradigm-shifting, AI-driven operating system architecture. '
    'Rather than relying on traditional, rigid OS scheduling routines that require human intervention, '
    'this system integrates an intelligent AI Agent at its core. The overriding goal is to build an environment '
    'where the OS acts as a proactive intelligence rather than a passive tool.'
)

doc.add_heading('2. The Dual-Partition Architecture', level=1)
doc.add_paragraph('Traditional OS architectures act as a single layer. Sunset applies a dual-partition approach:')
doc.add_heading('2.1. Normal OS Partition (User Environment)', level=2)
doc.add_paragraph('Houses standard applications (browsers, IDEs, games) and feels entirely unchanged to the end-user.')
doc.add_heading('2.2. AI Agent Partition (Intelligent Controller)', level=2)
doc.add_paragraph('An isolated, dedicated partition controlled exclusively by the AI. It continuously learns user patterns without privacy leaks and cannot be arbitrarily terminated.')
doc.add_heading('2.3. The Control Bridge (Secure API)', level=2)
doc.add_paragraph('A heavily monitored set of APIs connecting the Normal Partition to the AI Partition. Prevents unauthorized system access, guaranteeing high-grade cybersecurity.')

doc.add_heading('3. System Implementation & Process Logic', level=1)
doc.add_heading('Intelligent Memory Allocation', level=2)
doc.add_paragraph('Instead of generic swap-files, the AI Agent dynamically allocates RAM. It intervenes, saves states, and cleanly suspends background apps based on context.')
doc.add_heading('Modular AI Strategy', level=2)
doc.add_paragraph('Allows "updating the agent like an OS". Users can swap local inferences or offload heavy processing to safe cloud points while keeping telemetry on board.')

doc.add_heading('4. Target Applications', level=1)
doc.add_paragraph(
    '1. Smart System Optimization: Auto-detects slowness.\n'
    '2. Auto Task Manager: Contextually launches workflows.\n'
    '3. Privacy Control: Blocks erratic hardware access.\n'
    '4. Gaming Booster & Dev Assistant.\n'
    '5. Hybrid Processing capabilities.'
)

doc.save('d:/CodePlay/sunset-agent-os/Sunset_Agent_OS_Documentation.docx')
print("Successfully generated Word Doc.")

# 2. Generate PowerPoint Presentation
prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sunset Agent OS"
subtitle.text = "Next-Generation Dual-Partition AI Architecture\nAuthors: Arshad Pasha & Arun"

# Slide 2: Vision
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "The Core Vision"
tf = body.text_frame
tf.text = "Shift the traditional OS logic from static to intelligent."
p = tf.add_paragraph()
p.text = "Proactive intelligence manages memory boundaries natively."
p.level = 1
p = tf.add_paragraph()
p.text = "Seamlessly bridges a traditional OS with a dedicated AI controller."
p.level = 1

# Slide 3: Dual Partition
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Dual-Partition Architecture"
tf = body.text_frame
tf.text = "A reimagined system foundation dividing traditional interaction from AI control:"
p = tf.add_paragraph()
p.text = "1. Normal OS Partition (Standard User Environment)"
p.level = 1
p = tf.add_paragraph()
p.text = "2. AI Agent Partition (Intelligent Domain)"
p.level = 1
p = tf.add_paragraph()
p.text = "3. The Control Bridge (Secure communication via audited APIs)"
p.level = 1

# Slide 4: Resource Efficiency
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Intelligent Memory Resource Control"
tf = body.text_frame
tf.text = "How the system manages resources vs standard OS caches:"
p = tf.add_paragraph()
p.text = "Dynamic RAM allocation based on user patterns and current demands."
p.level = 1
p = tf.add_paragraph()
p.text = "Agent intervenes to cleanly state-save and throttle background process loops."
p.level = 1

# Slide 5: Modular Intelligence
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Modular System Evolution"
tf = body.text_frame
tf.text = "Updating the Brain:"
p = tf.add_paragraph()
p.text = "Replace / Update the Local LLM like a system driver."
p.level = 1
p = tf.add_paragraph()
p.text = "Keeps data secured strictly onboard inside the AI partition."
p.level = 1

# Slide 6: Ecosystem Benefits 
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Use Case Innovations"
tf = body.text_frame
tf.text = "What happens when an AI owns the System Controller?"
p = tf.add_paragraph()
p.text = "Zero-Configuration Game/Dev Boosting routines."
p.level = 1
p = tf.add_paragraph()
p.text = "Automated protection from malware scanning hardware behaviors natively."
p.level = 1

prs.save('d:/CodePlay/sunset-agent-os/Sunset_Agent_OS_Presentation.pptx')
print("Successfully generated PowerPoint.")
