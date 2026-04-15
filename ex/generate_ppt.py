from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
try:
    from PIL import Image
    import io
except ImportError:
    pass

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    if not os.path.exists(image_path):
        return False
    try:
        # Prevent format crashes by wrapping to standard PNG
        img = Image.open(image_path).convert('RGB')
        stream = io.BytesIO()
        img.save(stream, format='PNG')
        stream.seek(0)
        slide.shapes.add_picture(stream, left, top, width=width, height=height)
        return True
    except Exception as e:
        print(f"Skipping {image_path}: {e}")
        return False

def create_ppt():
    prs = Presentation()
    
    # ------------------ 1. TITLE SLIDE ------------------
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Smart Event Management System\n(SHRESHTA 2026)"
    title.text_frame.paragraphs[0].runs[0].font.bold = True
    subtitle.text = "A Real-Time Digital Transformation for Inter-Collegiate Fests\n\nPaper Presentation by:\n1. Mrs. Neha J K (Faculty Author)\n2. Arshad Pasha (Student Author)\nSeshadripuram Degree College, Mysuru"
    
    # ------------------ 2. INDEX ------------------
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "INDEX"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Abstract\n2. Introduction\n3. Literature Review\n4. Research Methodology\n5. Objectives of the Study\n6. Result and Discussion\n7. Conclusion\n8. References"
    
    # ------------------ 3. ABSTRACT ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "1. Abstract"
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "The Problem: Hosting large-scale inter-collegiate fests involves massive logistical hurdles—long queues, paper tracking, manual verification, and delayed certificates."
    p = tf.add_paragraph()
    p.text = "The Solution: We developed SHRESHTA to completely digitalize and streamline workflows autonomously."
    p = tf.add_paragraph()
    p.text = "Key Methodologies: Next.js client-side rendering, Firebase Firestore (NoSQL), HTML5 Canvas templates, and strict RBAC."
    p = tf.add_paragraph()
    p.text = "Impact: Achieved a 100% paperless cycle, reduced check-in time to under 5 seconds, and provided real-time transparency."

    # ------------------ 4. INTRODUCTION ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "2. Introduction"
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "The Current Context: Event management heavily relies on legacy workflows like Excel sheets and cash-based counters causing duplicate entries and tedious financial verification."
    p = tf.add_paragraph()
    p.text = "Our Digital Transformation: By leveraging modern web technologies and Firebase Serverless computing, we built a centralized application acting as a single source of truth."

    # ------------------ 5. LITERATURE REVIEW ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "3. Literature Review"
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "Traditional Methods: Simple forms provide a digital bridge but lack relational data constraints to block over-registration."
    p = tf.add_paragraph()
    p.text = "Commercial Architectures: Platforms like Eventbrite impose large fees and fail to accommodate academic needs like custom UTR tracking."
    p = tf.add_paragraph()
    p.text = "Smart Campuses: QR attendance proves barcode verification decreases queues by 80% (Kumar & Singh, 2022)."
    p = tf.add_paragraph()
    p.text = "The Gap: A severe lack of lightweight, real-time academic solutions."

    # ------------------ 6. RESEARCH METHODOLOGY ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.text = "4. Research Methodology & Architecture"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.text = "Built using Rapid Agile Development with a serverless macro-stack prioritizing real-time processing (Next.js 15, Firebase NoSQL, Tailwind CSS, Resend API)."
    add_image_safe(slide, "public/final_report/system_architecture.png", Inches(1.5), Inches(2.5), height=Inches(4.5))

    # ------------------ 7. OBJECTIVES ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "5. Objectives of the Study"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Eliminate Paper Trails: Store 100% of registrations on secured cloud."
    p = tf.add_paragraph()
    p.text = "2. Optimize Check-in Efficiency: Dramatically cut queue times."
    p = tf.add_paragraph()
    p.text = "3. Automate Validations: Establish a flawless digital payment (UTR) workflow."
    p = tf.add_paragraph()
    p.text = "4. Zero-Touch Rewards: Overlay names algorithmically for certificates."
    p = tf.add_paragraph()
    p.text = "5. Provide Real-Time Analytics via strict RBAC."

    # ------------------ 8. RESULTS OVERVIEW ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6. Result and Discussion (Overview)"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    tf = txBox.text_frame
    tf.text = "Live Deployment Success:"
    p = tf.add_paragraph()
    p.text = "The system successfully processed concurrent registrations for 13 distinct academic events."
    p = tf.add_paragraph()
    p.text = "Metric Precision: The real-time listener architecture completely obliterated manual counting—delivering live updates on Pendings and Check-ins."
    p = tf.add_paragraph()
    p.text = "0% Error Threshold: Isolated coordinator dashboards enforced absolute security, preventing data cross-contamination."
    add_image_safe(slide, "public/final_report/snapshot_home_page.png", Inches(2), Inches(4), height=Inches(3.2))

    # ------------------ 9. REGISTRATION ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6.1 Result: Registration Automation"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
    tf = txBox.text_frame
    tf.text = "Dynamic Form Intake:"
    p = tf.add_paragraph()
    p.text = "- UI dynamically mounts input rows based on event constraints."
    p = tf.add_paragraph()
    p.text = "- Captures 12-digit UTR immediately alongside a unified UPI QR."
    add_image_safe(slide, "public/final_report/reg_step3.png", Inches(4.5), Inches(1.5), height=Inches(4))

    # ------------------ 10. ADMIN DASHBOARDS ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6.2 Result: Real-Time Administration & RBAC"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    tf = txBox.text_frame
    tf.text = "Master & Coordinator Dashboards scale synchronously without browser refreshes."
    p = tf.add_paragraph()
    p.text = "Individual faculty coordinators view an identical dashboard that is mathematically restricted to display only their unique event participants."
    add_image_safe(slide, "public/final_report/snapshot_admin_dashboard.png", Inches(1), Inches(3), width=Inches(8))

    # ------------------ 11. CHECK IN DESK ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6.3 Result: Sub-5-Second Check-In Logistics"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.text = "Fuzzy Global Search queried immense rosters instantly. Physical check-in shrunk from 45 seconds to under 5 seconds."
    add_image_safe(slide, "public/final_report/cheeck_in_desk.png", Inches(2), Inches(3), width=Inches(6))

    # ------------------ 12. AUTOMATION & MAILS ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "6.4 Result: Zero-Touch Automation Flow"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = txBox.text_frame
    tf.text = "Instant API Dispatch via Resend mails QR tickets securely. Canvas generates HD participation certificates flawlessly without human oversight."
    add_image_safe(slide, "public/final_report/certificate_generaor.png", Inches(1.5), Inches(3), height=Inches(3))
    add_image_safe(slide, "public/final_report/email_recived.png", Inches(5), Inches(3), height=Inches(3))

    # ------------------ 13. CONCLUSION ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "7. Conclusion"
    tf = slide.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "The algorithm establishes definitively that high-friction physical event procedures can be entirely replaced by scalable micro-architectures."
    p = tf.add_paragraph()
    p.text = "Participants experienced frictionless onboarding."
    p = tf.add_paragraph()
    p.text = "Administrative faculty bypassed overlapping labor constraints."
    p = tf.add_paragraph()
    p.text = "The deployment stands as an infinitely reusable framework for future academic symposiums."

    # ------------------ 14. REFERENCES ------------------
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "8. References"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Vercel. (2026). Next.js Architecture Documentation."
    p = tf.add_paragraph()
    p.text = "2. Google. (2026). Cloud Firestore WebSocket Data Sync."
    p = tf.add_paragraph()
    p.text = "3. Sandhu, et al. (1996). Role-Based Access Control Models."
    p = tf.add_paragraph()
    p.text = "4. Kumar, R. (2022). QR Code Based Smart Attendance."
    p = tf.add_paragraph()
    p.text = "5. Resend. (2026). Serverless Email APIs."
    p = tf.add_paragraph()
    p.text = "6. MDN. (2026). Canvas API: Graphic Overlays."
    
    prs.save('SHRESHTA_Presentation_Slides_v2.pptx')
    print("PPTX generated successfully as SHRESHTA_Presentation_Slides_v2.pptx")

if __name__ == '__main__':
    create_ppt()
