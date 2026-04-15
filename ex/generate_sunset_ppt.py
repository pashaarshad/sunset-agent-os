from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
try:
    from PIL import Image
    import io
except ImportError:
    pass

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    if not os.path.exists(image_path):
        return False
    try:
        img = Image.open(image_path).convert('RGB')
        stream = io.BytesIO()
        img.save(stream, format='PNG')
        stream.seek(0)
        slide.shapes.add_picture(stream, left, top, width=width, height=height)
        return True
    except Exception as e:
        print(f"Skipping {image_path}: {e}")
        return False

def create_ppt():
    prs = Presentation()
    
    # ------------------ 1. TITLE SLIDE ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Sunset Agent OS"
    subtitle.text = "A Paradigm-Shifting AI-Driven OS Architecture\n\nAuthors:\nArshad Pasha (Student Author)\nArun (Research Lead)\n\nSeshadripuram Degree College"

    # ------------------ 2. INDEX ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "INDEX"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Abstract & Vision\n2. The Problem: Legacy OS Limits\n3. Core Concept: Dual Partition System\n4. Control Bridge & Security\n5. Modular Intelligence (LLM Switching)\n6. Real-World Applications\n7. Future Roadmap\n8. Conclusion"

    # ------------------ 3. ABSTRACT & VISION ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Abstract & Vision"
    tf = slide.placeholders[1].text_frame
    tf.text = "Vision: Shift from Static OS logic to Proactive Intelligence."
    p = tf.add_paragraph()
    p.text = "Core: AI Agent becomes the core decision-maker."
    p = tf.add_paragraph()
    p.text = "Model: User → AI Agent → OS → Hardware."

    # ------------------ 4. THE PROBLEM ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. The Problem: Legacy OS Limits"
    tf = slide.placeholders[1].text_frame
    tf.text = "- Reactive Resource Management (user closes memory-heavy apps)."
    p = tf.add_paragraph()
    p.text = "- Static Schedulers (no context of user behavior)."
    p = tf.add_paragraph()
    p.text = "- Privacy Risks (cloud-based processing in current AI solutions)."
    p = tf.add_paragraph()
    p.text = "- Fixed Intelligence (cannot swap brains without reinstalling)."

    # ------------------ 5. CORE CONCEPT: DUAL PARTITION ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "3. Core Concept: Dual Partition System"
    add_image_safe(slide, "architecture/dual_partition_diagram.png", Inches(2), Inches(1.5), height=Inches(5))

    # ------------------ 6. PARTITION BREAKDOWN ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Partitions: Owner vs. User"
    tf = slide.placeholders[1].text_frame
    tf.text = "AI Agent Partition (Agent Owned):"
    p = tf.add_paragraph()
    p.text = "- Isolated, Secure environment for Analysis."
    p = tf.add_paragraph()
    p.text = "- Cannot be terminated by user-side processes."
    p = tf.add_paragraph()
    p.text = "\nNormal OS Partition (User Shared):"
    p = tf.add_paragraph()
    p.text = "- Familiar Windows/Linux Apps."
    p = tf.add_paragraph()
    p.text = "- Data isolation via secure Control Bridge."

    # ------------------ 7. CONTROL BRIDGE & SECURITY ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "4. Control Bridge & Local Security"
    add_image_safe(slide, "architecture/new_system_architecture.png", Inches(1), Inches(1.5), width=Inches(8))
    # Add small textbox about security
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
    tx.text_frame.text = "Security: Data stays local. AI processes behavior in its own domain. No cloud leaks."

    # ------------------ 8. MODULAR INTELLIGENCE ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. Modular Intelligence (LLM Switching)"
    tf = slide.placeholders[1].text_frame
    tf.text = "- AI = Modular Brain of the OS."
    p = tf.add_paragraph()
    p.text = "- Swap Models: Change LLM (GPT/Llama/Local) like a plugin."
    p = tf.add_paragraph()
    p.text = "- Upgradable: Update Agent independently from OS Kernel."
    p = tf.add_paragraph()
    p.text = "- Performance: Choose brain size based on hardware capability."

    # ------------------ 9. APPLICATIONS 1 ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Real-World Applications (Efficiency)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Smart System Optimization: Auto-RAM clear and CPU throttling."
    p = tf.add_paragraph()
    p.text = "Auto Task Manager: Context-aware app launching."
    p = tf.add_paragraph()
    p.text = "Gaming Booster: Directs max performance to active game window."
    p = tf.add_paragraph()
    p.text = "Dev Assistant: Environment setup and dependency management."

    # ------------------ 10. APPLICATIONS 2 ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Real-World Applications (Security)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Behavior Analysis: Predicts user needs and suggests improvements."
    p = tf.add_paragraph()
    p.text = "Privacy Control: Hardware-level blocks for suspicious apps."
    p = tf.add_paragraph()
    p.text = "Hybrid Processing: Sensitive tasks stay local; heavy tasks route safely."

    # ------------------ 11. AGENT FLOW ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "AI Agent Process Flow"
    add_image_safe(slide, "images/ai_agent_flow.png", Inches(1), Inches(1.5), width=Inches(8))

    # ------------------ 12. FUTURE VISION ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "7. Future Roadmap & Vision"
    add_image_safe(slide, "images/future_vision.png", Inches(3), Inches(1.5), height=Inches(5))
    tx = slide.shapes.add_textbox(Inches(0.2), Inches(2), Inches(3), Inches(3))
    tx.text_frame.text = "Steps:\n1. Python Prototype\n2. Bridge API\n3. Linux Integration\n4. Hardware Partition"

    # ------------------ 13. CONCLUSION ------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "8. Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Sunset Agent OS is THE future of computing."
    p = tf.add_paragraph()
    p.text = "Intelligence + Privacy + Performance."
    p = tf.add_paragraph()
    p.text = "The system is no longer just a tool—it is an intelligent agent."
    
    save_path = 'Sunset_Agent_OS_Presentation_Final.pptx'
    prs.save(save_path)
    print(f"PPT generated successfully: {save_path}")

if __name__ == '__main__':
    create_ppt()
